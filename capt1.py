from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
#slide 1
slide1_register= prs.slide_layouts[1]
slide1=prs.slides.add_slide(slide1_register)

#Edit Bullet Point slide
title1=slide1.shapes.title
title1.text="OCS"

bullet_point_box=slide1.shapes
bullet_points_lvl1=bullet_point_box.placeholders[1]
bullet_points_lvl1.text="CAPS"
#Add Image
img1="screen/0.png"

left = top = Inches(2.5)


left = Inches(0.5)
height = Inches(4)
pic = slide1.shapes.add_picture(img1, left, top, height=height)

################################################################################################
#slide2
#Create bullet Point slide
slide2_register= prs.slide_layouts[1]
slide2=prs.slides.add_slide(slide2_register)

#Edit Bullet Point slide
title2=slide2.shapes.title
title2.text="OCS"

bullet_point_box=slide2.shapes
bullet_points_lvl1=bullet_point_box.placeholders[1]
bullet_points_lvl1.text="Diameter Data Counters (Data Sessions)"
#Add Image
img1="screen/1.png"

left = top = Inches(2.5)


left = Inches(0.5)
height = Inches(4)
pic = slide2.shapes.add_picture(img1, left, top, height=height)
##########################################################################################################
slide3_register= prs.slide_layouts[1]
slide3=prs.slides.add_slide(slide3_register)

#Edit Bullet Point slide
title3=slide3.shapes.title
title3.text="OCS"

bullet_point_box=slide3.shapes
bullet_points_lvl1=bullet_point_box.placeholders[1]
bullet_points_lvl1.text="Diameter Open data (Data Sessions)"
#Add Image
img1="screen/2.png"

left = top = Inches(2.5)


left = Inches(0.5)
height = Inches(4)
pic = slide3.shapes.add_picture(img1, left, top, height=height)

##########################################################################################################
slide4_register= prs.slide_layouts[1]
slide4=prs.slides.add_slide(slide4_register)

#Edit Bullet Point slide
title4=slide4.shapes.title
title4.text="OCS"

bullet_point_box=slide4.shapes
bullet_points_lvl1=bullet_point_box.placeholders[1]
bullet_points_lvl1.text="Diameter active Sessions"
#Add Image
img1="screen/3.png"

left = top = Inches(2.5)


left = Inches(0.5)
height = Inches(4)
pic = slide3.shapes.add_picture(img1, left, top, height=height)

##########################################################################################################
slide5_register= prs.slide_layouts[1]
slide5=prs.slides.add_slide(slide5_register)

#Edit Bullet Point slide
title5=slide4.shapes.title
title5.text="OCS"

bullet_point_box=slide5.shapes
bullet_points_lvl1=bullet_point_box.placeholders[1]
bullet_points_lvl1.text="Diameter Call Active (Le nombre de call)"
#Add Image
img1="screen/4.png"

left = top = Inches(2.5)


left = Inches(0.5)
height = Inches(4)
pic = slide4.shapes.add_picture(img1, left, top, height=height)
#enrigistrement de la page 
prs.save('test.pptx')
